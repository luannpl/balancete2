from PySide6 import QtWidgets, QtGui, QtCore
from PySide6.QtWidgets import QMessageBox, QTableWidgetItem, QFileDialog, QInputDialog, QDialog, QVBoxLayout, QLabel, QLineEdit, QListWidget, QPushButton, QComboBox
from PySide6.QtGui import QKeySequence, QShortcut
from PySide6.QtCore import Qt
from utils.icone import usar_icone
from utils.mensagem import mensagem_error, mensagem_sucesso, mensagem_aviso
from db.conexao import conectar_com_banco, fechar_banco, criar_tabela, tabela_slide, codigos_slide, tabela_empresa, tabela_usuario
from middleware.processamento import processar_arquivo
import sys
from datetime import datetime
import os
import requests
import mysql.connector
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
# from pptx.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

class UserLogin(QtWidgets.QWidget):
    def __init__(self):
        super(UserLogin, self).__init__()
        self.setWindowTitle('Balancete - Login')
        self.setGeometry(200, 200, 900, 700)
        self.setStyleSheet('background-color: #030d18;')

        self.layout = QtWidgets.QVBoxLayout()
        self.layout.setContentsMargins(20, 20, 20, 20)
        self.layout.setSpacing(10)

        self.title_label = QtWidgets.QLabel('Login - Balancete')
        self.title_label.setAlignment(QtCore.Qt.AlignLeft)
        self.title_label.setStyleSheet("font-size: 36px; font-weight: bold; color: #001F3F; margin: 0px; padding: 0px;")

        self.user_label = QtWidgets.QLabel('Usuário:')
        self.user_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.user_input = QtWidgets.QLineEdit()
        self.user_input.setPlaceholderText('Digite o usuário')
        self.user_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;") 

        self.password_label = QtWidgets.QLabel('Senha:')
        self.password_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.password_input = QtWidgets.QLineEdit()
        self.password_input.setEchoMode(QtWidgets.QLineEdit.Password)
        self.password_input.setPlaceholderText('Digite a senha')
        self.password_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;")

        self.login_button = QtWidgets.QPushButton('Login')
        self.login_button.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.login_button.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.login_button.clicked.connect(self.validate_login)

        self.cadastrar_button = QtWidgets.QPushButton('Cadastrar')
        self.cadastrar_button.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.cadastrar_button.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.cadastrar_button.clicked.connect(self.cadastrar_user)

        group_box = QtWidgets.QGroupBox()
        group_box.setStyleSheet("background-color: #FFFFFF; padding: 20px; border-radius: 10px;")
        group_box.setMaximumWidth(500) 
        group_box.setMinimumWidth(500)  

        group_box_layout = QtWidgets.QVBoxLayout()
        group_box_layout.setSpacing(10)  
        group_box_layout.setContentsMargins(10, 10, 10, 10)
        group_box_layout.addWidget(self.title_label) 
        group_box_layout.addWidget(self.user_label)
        group_box_layout.addWidget(self.user_input)
        group_box_layout.addWidget(self.password_label)
        group_box_layout.addWidget(self.password_input)
        group_box_layout.addWidget(self.login_button)
        group_box_layout.addWidget(self.cadastrar_button)
        group_box.setLayout(group_box_layout)

        center_layout = QtWidgets.QVBoxLayout()
        center_layout.addStretch()
        center_layout.addWidget(group_box, alignment=QtCore.Qt.AlignCenter)
        center_layout.addStretch()

        main_layout = QtWidgets.QVBoxLayout()
        main_layout.addLayout(center_layout)

        self.setLayout(main_layout)
    
    def validate_login(self):
        user = self.user_input.text().strip()
        password = self.password_input.text().strip()
        conexao = conectar_com_banco()
        tabela_usuario(conexao)
        cursor = conexao.cursor()
        cursor.execute("USE railway")
        cursor.execute("SELECT senha FROM user WHERE nome = %s", (user,))
        result = cursor.fetchone()
        
        if not result:
            mensagem_error("Usuário ou senha inválidos.")
            return
        result = result[0]
        

        if password == result:
            self.empresa_window = EmpresaWindow(user)
            usar_icone(self.empresa_window)
            self.empresa_window.showMaximized()
            self.close()
        else:
            mensagem_error("Usuário ou senha inválidos.")
            return
    
    def cadastrar_user(self):
        self.cadastro = UserCadastro()
        usar_icone(self.cadastro)
        self.cadastro.showMaximized()
        self.close()

class UserCadastro(QtWidgets.QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle('Balancete - Cadastro')
        self.setGeometry(200, 200, 900, 700)
        self.setStyleSheet('background-color: #030d18;')

        self.layout = QtWidgets.QVBoxLayout()
        self.layout.setContentsMargins(20, 20, 20, 20)
        self.layout.setSpacing(10)

        self.title_label = QtWidgets.QLabel('Cadastro - Balancete')
        self.title_label.setAlignment(QtCore.Qt.AlignLeft)
        self.title_label.setStyleSheet("font-size: 36px; font-weight: bold; color: #001F3F; margin: 0px; padding: 0px;")

        self.user_label = QtWidgets.QLabel('Usuário:')
        self.user_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.user_input = QtWidgets.QLineEdit()
        self.user_input.setPlaceholderText('Digite o usuário')
        self.user_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;") 

        # Campo de entrada para senha
        self.password_label = QtWidgets.QLabel('Senha:')
        self.password_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.password_input = QtWidgets.QLineEdit()
        self.password_input.setEchoMode(QtWidgets.QLineEdit.Password)
        self.password_input.setPlaceholderText('Digite a senha')
        self.password_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;")

        self.password_label_confirm = QtWidgets.QLabel('Confirme a senha:')
        self.password_label_confirm.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.password_input_confirm = QtWidgets.QLineEdit()
        self.password_input_confirm.setEchoMode(QtWidgets.QLineEdit.Password)
        self.password_input_confirm.setPlaceholderText('Confirme a senha')
        self.password_input_confirm.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;")

        # Botão para login
        self.cadastrar_button = QtWidgets.QPushButton('Cadastrar')
        self.cadastrar_button.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.cadastrar_button.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.cadastrar_button.clicked.connect(self.cadastrar_user)

        self.voltar_btn = QtWidgets.QPushButton('Login')
        self.voltar_btn.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.voltar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.voltar_btn.clicked.connect(self.voltar)

        group_box = QtWidgets.QGroupBox()
        group_box.setStyleSheet("background-color: #FFFFFF; padding: 20px; border-radius: 10px;")
        group_box.setMaximumWidth(500)  # Aumentando a largura máxima do grupo
        group_box.setMinimumWidth(500)  # Define uma largura mínima para evitar colapso

        # Adiciona o layout ao grupo
        group_box_layout = QtWidgets.QVBoxLayout()
        group_box_layout.setSpacing(10)  # Ajusta o espaçamento vertical entre os widgets
        group_box_layout.setContentsMargins(10, 10, 10, 10)  # Ajusta margens internas do grupo
        group_box_layout.addWidget(self.title_label)
        group_box_layout.addWidget(self.user_label)
        group_box_layout.addWidget(self.user_input)
        group_box_layout.addWidget(self.password_label)
        group_box_layout.addWidget(self.password_input)
        group_box_layout.addWidget(self.password_label_confirm)
        group_box_layout.addWidget(self.password_input_confirm)
        group_box_layout.addWidget(self.cadastrar_button)
        group_box_layout.addWidget(self.voltar_btn)
        group_box.setLayout(group_box_layout)

        center_layout = QtWidgets.QVBoxLayout()
        center_layout.addStretch()
        center_layout.addWidget(group_box, alignment=QtCore.Qt.AlignCenter)
        center_layout.addStretch()

        main_layout = QtWidgets.QVBoxLayout()
        main_layout.addLayout(center_layout)

        self.setLayout(main_layout)    

    def cadastrar_user(self):
        user = self.user_input.text().strip()
        password = self.password_input.text().strip()
        password_confirm = self.password_input_confirm.text().strip()
        print(user, password)
        print(len(user), len(password))
        if len(user) < 1:
            mensagem_aviso("Usuário precisa ter pelo menos um caractere.")
            return
        if len(password) < 1:
            mensagem_aviso("Senha precisa ter pelo menos um caractere.")
            return

        if password != password_confirm:
            mensagem_error("Senhas não conferem.")
            return

        conexao = conectar_com_banco()
        tabela_usuario(conexao)
        cursor = conexao.cursor()
        cursor.execute("USE railway")
        cursor.execute("INSERT INTO user (nome, senha) VALUES (%s, %s)", (user, password))
        conexao.commit()
        mensagem_sucesso("Usuário cadastrado com sucesso.")

    def voltar(self):
        self.login = UserLogin()
        usar_icone(self.login)
        self.login.showMaximized()
        self.close()    

class EmpresaWindow(QtWidgets.QWidget):
    def __init__(self, user):
        super().__init__()
        self.setWindowTitle('Balancete - Empresas')
        self.setGeometry(200, 200, 900, 700)
        self.setStyleSheet('background-color: #030d18;')

        self.user = user
        self.layout = QtWidgets.QVBoxLayout()
        self.layout.setContentsMargins(20, 20, 20, 20)
        self.layout.setSpacing(10)

        self.label_titulo = QtWidgets.QLabel('Escolha ou cadastre uma empresa')
        self.label_titulo.setStyleSheet("font-size: 20px; font-weight: bold; color: #ffffff; margin: 0px; padding: 0px;")
        self.label_titulo.setAlignment(QtCore.Qt.AlignCenter)
        self.layout.addStretch()
        self.layout.addWidget(self.label_titulo)

        self.layout.addSpacing(20)

        self.combo_empresas = QtWidgets.QComboBox()
        self.combo_empresas.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color:rgb(255, 255, 255);")
        self.combo_empresas.setFixedWidth(400)
        self.carregar_empresas()
    
        self.layout.addWidget(self.combo_empresas, alignment=QtCore.Qt.AlignCenter)

        self.layout.addSpacing(20)

        self.entrar_btn = QtWidgets.QPushButton('Entrar')
        self.entrar_btn.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.entrar_btn.setFixedWidth(400)
        self.entrar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.entrar_btn.clicked.connect(self.entrar)

        self.layout.addWidget(self.entrar_btn, alignment=QtCore.Qt.AlignCenter)

        self.layout.addSpacing(10)
        self.cadastrar_btn = QtWidgets.QPushButton('Cadastrar Empresa')
        self.cadastrar_btn.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.cadastrar_btn.setFixedWidth(400)
        self.cadastrar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.cadastrar_btn.clicked.connect(self.cadastrar_empresa)

        self.layout.addWidget(self.cadastrar_btn, alignment=QtCore.Qt.AlignCenter)

        self.layout.addStretch()
        self.setLayout(self.layout)

    def entrar(self):
        empresa = self.combo_empresas.currentText()
        if empresa == "Selecione uma empresa":
            mensagem_aviso("Selecione uma empresa.")
            return
        self.janela_principal = MainWindow(empresa, self.user)
        usar_icone(self.janela_principal)
        self.janela_principal.showMaximized()
        self.close()
    
    def cadastrar_empresa(self):
        self.empresa_cadastro = EmpresaCadastro(self.user)
        usar_icone(self.empresa_cadastro)
        self.empresa_cadastro.showMaximized()
        self.close()

    def carregar_empresas(self):
        conexao = conectar_com_banco()
        tabela_empresa(conexao)
        try:
            cursor = conexao.cursor()
            cursor.execute("USE railway")
            cursor.execute("SELECT razao_social FROM empresas ORDER BY razao_social ASC")
            empresas = [row[0] for row in cursor.fetchall()]
            
            self.combo_empresas.clear()
            self.combo_empresas.addItem("Selecione uma empresa")
            self.combo_empresas.addItems(empresas)
            self.combo_empresas.model().item(0).setEnabled(False)

        except Exception as e:
            QtWidgets.QMessageBox.critical(self, "Erro", f"Ocorreu um erro ao acessar o banco de dados: {e}")
        finally:
            if conexao.is_connected():
                cursor.close()

class EmpresaCadastro(QtWidgets.QWidget):
    def __init__(self, user):
        super().__init__()
        self.setWindowTitle('Balancete - Cadastro de Empresas')
        self.setGeometry(200, 200, 900, 700)
        self.setStyleSheet('background-color: #030d18;')

        self.user = user

        self.layout = QtWidgets.QVBoxLayout()
        self.layout.setContentsMargins(20, 20, 20, 20)
        self.layout.setSpacing(10)

        self.title_label = QtWidgets.QLabel('Cadastro de Empresa')
        self.title_label.setAlignment(QtCore.Qt.AlignLeft)
        self.title_label.setStyleSheet("font-size: 36px; font-weight: bold; color: #001F3F; margin: 0px; padding: 0px;")

        self.cnpj_label = QtWidgets.QLabel('CNPJ:')
        self.cnpj_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.cnpj_input = QtWidgets.QLineEdit()
        self.cnpj_input.setInputMask('99.999.999/9999-99')
        self.cnpj_input.setPlaceholderText('Digite o CNPJ')
        self.cnpj_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;")

        self.razao_social_label = QtWidgets.QLabel('Razão Social:')
        self.razao_social_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #000000; margin: 0px; padding: 0px;")
        self.razao_social_input = QtWidgets.QLineEdit()
        self.razao_social_input.setPlaceholderText('Digite a razão social')
        self.razao_social_input.setStyleSheet("font-size: 20px; padding: 8px; border: 1px solid #ccc; border-radius: 5px; color: #000000;")

        self.btn_cadastrar_empresa = QtWidgets.QPushButton('Cadastrar')
        self.btn_cadastrar_empresa.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.btn_cadastrar_empresa.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.btn_cadastrar_empresa.clicked.connect(self.cadastrar_empresa)

        self.voltar_btn = QtWidgets.QPushButton('Voltar')
        self.voltar_btn.setStyleSheet("""QPushButton {
            font-size: 20px;
            font-weight: bold;
            padding: 10px;
            background-color: #001F3F;
            color: #ffffff;
        }
        QPushButton:hover {
            background-color: #005588;  /* Cor de fundo quando o mouse passa sobre o botão */
            color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
        }""")
        self.voltar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.voltar_btn.clicked.connect(self.voltar)

        group_box = QtWidgets.QGroupBox()
        group_box.setStyleSheet("background-color: #FFFFFF; padding: 20px; border-radius: 10px;")
        group_box.setMaximumWidth(500)
        group_box.setMinimumWidth(500)

        group_box_layout = QtWidgets.QVBoxLayout()
        group_box_layout.setSpacing(10)
        group_box_layout.setContentsMargins(10, 10, 10, 10)
        group_box_layout.addWidget(self.title_label)
        group_box_layout.addWidget(self.cnpj_label)
        group_box_layout.addWidget(self.cnpj_input)
        group_box_layout.addWidget(self.razao_social_label)
        group_box_layout.addWidget(self.razao_social_input)
        group_box_layout.addWidget(self.btn_cadastrar_empresa)
        group_box_layout.addWidget(self.voltar_btn)
        group_box.setLayout(group_box_layout)

        center_layout = QtWidgets.QVBoxLayout()
        center_layout.addStretch()
        center_layout.addWidget(group_box, alignment=QtCore.Qt.AlignCenter)
        center_layout.addStretch()

        main_layout = QtWidgets.QVBoxLayout()
        main_layout.addLayout(center_layout)

        self.setLayout(main_layout)
    
    def cadastrar_empresa(self):
        cnpj = self.cnpj_input.text().strip()
        razao_social = self.razao_social_input.text().strip()

        if len(cnpj) != 18:
            mensagem_error("CNPJ inválido.")
            return
        
        if len(razao_social) < 1:
            mensagem_error("Razão social inválida.")
            return
        
        try:
            conexao = conectar_com_banco()
            tabela_empresa(conexao)
            cursor = conexao.cursor()
            cursor.execute("USE railway")
            cursor.execute("INSERT INTO empresas (cnpj, razao_social, user) VALUES (%s, %s, %s)", (cnpj, razao_social, self.user))
            conexao.commit()
            mensagem_sucesso("Empresa cadastrada com sucesso.")
        except Exception as e:
            mensagem_error(f"Ocorreu um erro ao cadastrar a empresa: {e}")
        finally:
            if conexao.is_connected():
                cursor.close()
                conexao.close()

    
    def voltar(self):
        self.empresas = EmpresaWindow(self.user)
        usar_icone(self.empresas)
        self.empresas.showMaximized()
        self.close()

class MainWindow(QtWidgets.QMainWindow):
    def __init__(self, empresa, user):
        super().__init__()
        self.empresa = empresa
        self.user = user

        self.setWindowTitle('Balancete - Dashboard')
        self.setGeometry(200, 200, 900, 700)
        self.setStyleSheet('background-color: #030d18;')

        self.central_widget = QtWidgets.QWidget()
        self.setCentralWidget(self.central_widget)
        self.layout = QtWidgets.QVBoxLayout(self.central_widget)

        self.header = QtWidgets.QGroupBox('Empresa')
        self.frame_layout = QtWidgets.QVBoxLayout(self.header)

        self.empresa_label = QtWidgets.QLabel(f'Empresa: {self.empresa}')
        self.empresa_label.setStyleSheet("font-size: 20px; font-weight: bold; color: #ffffff; margin: 0px; padding: 0px;")
        self.frame_layout.addWidget(self.empresa_label)

        self.layout.addWidget(self.header)

        self.btn_frame = QtWidgets.QHBoxLayout()
        self.layout.addLayout(self.btn_frame)

        self.criar_botoes()
        self.label_arquivo = QtWidgets.QLabel("Nenhum arquivo selecionado.")
        self.layout.addWidget(self.label_arquivo)
        self.criar_combo_boxes()
        self.criar_tabela()

        self.stack_layout = QtWidgets.QStackedLayout()
        self.layout.addLayout(self.stack_layout)

        self.combo_box_slides.currentIndexChanged.connect(lambda: self.gerar_tabela_se_valido(self.tabela, self.combo_box_slides.currentText(), self.empresa, self.combo_ano.currentText(), self.stack_layout))
        self.combo_box_meses.currentIndexChanged.connect(lambda: self.gerar_tabela_se_valido(self.tabela, self.combo_box_meses.currentText(), self.empresa, self.combo_ano.currentText(), self.stack_layout))

        self.imagem_placeholder = QtWidgets.QLabel()
        self.imagem_placeholder.setPixmap(QtGui.QPixmap("images\\logo.png").scaled(300, 300, QtCore.Qt.KeepAspectRatio))
        self.imagem_placeholder.setAlignment(QtCore.Qt.AlignCenter)

        self.stack_layout.addWidget(self.imagem_placeholder)
        self.stack_layout.addWidget(self.tabela)

        self.stack_layout.setCurrentWidget(self.imagem_placeholder)

        self.url_icone = "https://cdn-icons-png.flaticon.com/512/4726/4726016.png"

        self.caminho_icone_local = os.path.join(os.path.dirname(__file__), "images", "icon.png")

        if not os.path.exists(self.caminho_icone_local):
            self.baixar_icone(self.url_icone, self.caminho_icone_local)
            
        self.caminho_logo = self.recurso_caminho("images\\logo.png")
        self.botao_exportar = QtWidgets.QPushButton()
        self.botao_exportar.setIcon(QtGui.QIcon(self.caminho_icone_local))
        self.botao_exportar.setIconSize(QtCore.QSize(32, 32))
        self.botao_exportar.setStyleSheet("""
            QPushButton {
                background-color: transparent;
                border: none;
            }
            QPushButton:hover {
                background-color: #FFA500;
            }
        """)
        self.botao_exportar.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        self.botao_exportar.clicked.connect(lambda: self.exportar_tabela_para_powerpoint(
        self.tabela,
        self.combo_box_meses.currentText() if self.combo_box_meses.currentText() != "Balancetes" else self.combo_box_slides.currentText(), 
        self.empresa,
        self.caminho_logo  
        ))

        self.layout_botoes = QtWidgets.QHBoxLayout()
        self.layout_botoes.addWidget(self.botao_exportar, alignment=QtCore.Qt.AlignLeft | QtCore.Qt.AlignTop)

        self.layout.addLayout(self.layout_botoes)

        self.progress_bar = QtWidgets.QProgressBar()
        self.layout.addWidget(self.progress_bar)

        self.status_bar = QtWidgets.QStatusBar()
        self.setStatusBar(self.status_bar)

        self.status_bar.showMessage("Assertivus Contábil - Balancete")
    


        self.setLayout(self.layout)
    
    
    def inserir_dados_mensal(self, progress_bar, label_arquivo,  empresa, nome_usuario, stack_layout):
        stack_layout.setCurrentWidget(self.imagem_placeholder)
        balancete_mensal = 'balancete_mensal'
        arquivo, _ = QFileDialog.getOpenFileName(self, "Selecione o Arquivo", "", "Arquivos Excel (*.xlsx)")

        if not arquivo:
            label_arquivo.setText("Nenhum arquivo selecionado.")
            return

        label_arquivo.setText(f"Arquivo Selecionado: {arquivo}")
        progress_bar.setValue(10)

        dados = processar_arquivo(arquivo)
        if dados is None:
            mensagem_error("Erro ao processar o arquivo!")
            progress_bar.setValue(0)
            return

        progress_bar.setValue(40)

        meses = [
            "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho",
            "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro"
        ]
        ano_atual = datetime.now().year
        anos = [str(ano_atual-1), str(ano_atual), str(ano_atual+1), str(ano_atual+2)]  # Exemplo: anos de 2020 até 2030
        mes_selecionado, ok_mes = QInputDialog.getItem(
            self,
            "Selecione o Mês",
            "Escolha o mês para associar ao arquivo:",
            meses,
            0,
            False
        )
        

        if not ok_mes:
            mensagem_aviso("Nenhum mês foi selecionado.")
            progress_bar.setValue(0)
            return

        ano_selecionado, ok_ano = QInputDialog.getItem(
            self,
            "Selecione o Ano",
            "Escolha o ano para associar ao arquivo:",
            anos,
            1,
            False
        )

        if not ok_ano:
            mensagem_aviso("Nenhum ano foi selecionado.")
            progress_bar.setValue(0)
            return

        mes_numero = meses.index(mes_selecionado) + 1
        mes_ano = f"{mes_numero:02d}-{ano_selecionado}"  # Formato: MM-AAAA, exemplo: "01-2024"

        progress_bar.setValue(60)

        conexao = conectar_com_banco()
        if conexao:
            # Verifica se já existe dado para o período selecionado
            if self.verificar_periodo_existente(conexao, balancete_mensal, mes_ano, empresa):
                confirmacao = QtWidgets.QMessageBox.question(
                    self,
                    "Confirmação",
                    f"Já existem dados para o período de {mes_selecionado}/{ano_selecionado}. Deseja continuar e sobrescrever os dados?",
                    QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No,
                )
                if confirmacao == QtWidgets.QMessageBox.No:
                    progress_bar.setValue(0)
                    fechar_banco(conexao)
                    return

            criar_tabela(conexao)
            progress_bar.setValue(70)

            dados_para_inserir = [
                (
                    row['Conta'],
                    row.get('Descrição', ''),
                    row.get('Saldo Anterior', ''),
                    row.get('Débitos', ''),
                    row.get('Créditos', ''),
                    row.get('Saldo Atual', ''),
                    nome_usuario,
                    mes_ano  # Usando mes_ano no lugar de mes_numero
                )
                for _, row in dados.iterrows()
            ]

            self.insert_data(conexao, dados_para_inserir, balancete_mensal, mes_ano)
            progress_bar.setValue(90)

            fechar_banco(conexao)
            progress_bar.setValue(100)
            mensagem_sucesso("Dados inseridos com sucesso!")
            progress_bar.setValue(0)
        else:
            mensagem_error("Erro ao conectar ao banco de dados.")
            progress_bar.setValue(0)

    def verificar_periodo_existente(self, conexao, tabela, periodo, empresa):
        """
        Verifica se já existem dados para o período no banco.
        """
        try:
            cursor = conexao.cursor()
            query = f"SELECT COUNT(*) FROM {tabela} WHERE periodo = %s AND empresa = %s"
            cursor.execute(query, (periodo, empresa))
            resultado = cursor.fetchone()
            return resultado[0] > 0  # Retorna True se existir pelo menos um registro
        except mysql.connector.Error as e:
            print(f"Erro ao verificar período existente: {e}")
            return False
        finally:
            cursor.close()
    
    def insert_data(self,conexao, data, balancete, mes_ano):
        try:
            cursor = conexao.cursor()

            insert_query = f"""
            INSERT INTO {balancete} (Conta, Descrição, `Saldo Anterior`, Débitos, Créditos, `Saldo Atual`, Usuario, empresa, periodo)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s);
            """

            dados_com_usuario = [(row[0], row[1], row[2], row[3], row[4], row[5], self.user, self.empresa, mes_ano) for row in data]

            cursor.executemany(insert_query, dados_com_usuario)
            conexao.commit()
            # print("Dados inseridos com sucesso!")
                
        except mysql.connector.Error as e:
            mensagem_error(f"Erro ao inserir dados: {e}")
        finally:
            cursor.close()

    def criar_botoes(self):
        self.botoes_info = [
            ("Inserir Balancete", lambda: self.inserir_dados_mensal(self.progress_bar, self.label_arquivo, self.empresa, self.user, self.stack_layout)),
            ("Criar Slide", lambda: self.criar_slide(self.combo_box_slides, self.tabela, self.empresa, self.stack_layout)),
            ("Editar Slide", lambda: self.editar_slide(self.combo_box_slides, self.tabela, self.empresa, self.stack_layout)),
        ]

        for texto, comando in self.botoes_info:
            botao = QtWidgets.QPushButton(texto)
            botao.clicked.connect(lambda: self.reset_combo_slides())
            botao.clicked.connect(lambda: self.reset_combo())
            botao.clicked.connect(lambda: self.limpar_tabela(self.tabela))
            botao.setStyleSheet("""QPushButton {
                font-size: 20px;
                font-weight: bold;
                padding: 10px;
                background-color: #001F3F;
                color: #ffffff;
            }
            QPushButton:hover {
                background-color: #005588 !important;  /* Cor de fundo quando o mouse passa sobre o botão */
                color:#ffffff;  /* Cor do texto quando o mouse passa sobre o botão */
            }""")
            botao.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
            botao.clicked.connect(comando)
            self.btn_frame.addWidget(botao)
    
    def criar_combo_boxes(self):
        self.combo_box_layout = QtWidgets.QHBoxLayout()
        
        ano_atual = datetime.now().year
        self.combo_ano = QtWidgets.QComboBox()
        self.combo_ano.addItems(["Escolha o ano", str(ano_atual - 1), str(ano_atual), str(ano_atual + 1), str(ano_atual + 2)])
        self.combo_ano.model().item(0).setEnabled(False)
        self.combo_ano.setStyleSheet("font-size: 16px; padding: 5px; color: #000000; background-color: #b0e0e6;")
        
        self.combo_box_slides = QtWidgets.QComboBox()
        self.combo_box_slides.setStyleSheet("font-size: 16px; padding: 5px; color: #000000; background-color: #b0e0e6;")
        self.combo_box_meses = QtWidgets.QComboBox()
        self.combo_box_meses.setStyleSheet("font-size: 16px; padding: 5px; color: #000000; background-color: #b0e0e6;")
        
        conexao = conectar_com_banco()
        tabela_slide(conexao)
        codigos_slide(conexao)
        cursor = conexao.cursor()
        
        cursor.execute("SELECT id FROM empresas WHERE razao_social = %s", (self.empresa,))
        id_empresa = cursor.fetchone()[0]
        
        cursor.execute("SELECT nome FROM slides WHERE id_empresa = %s", (id_empresa,))
        itens_slides = ["Slides"] + [linha[0] for linha in cursor.fetchall()]
        self.combo_box_slides.addItems(itens_slides)
        
        itens_meses = ["Balancetes", "Janeiro", "Fevereiro", "Março", "Abril", "Maio", "Junho", "Julho", "Agosto", "Setembro", "Outubro", "Novembro", "Dezembro", "Balancete Acumulado"]
        self.combo_box_meses.addItems(itens_meses)

        self.combo_box_slides.model().item(0).setEnabled(False)
        self.combo_box_meses.model().item(0).setEnabled(False)

        self.combo_ano.activated.connect(lambda: self.limpar_tabela(self.tabela))
        self.combo_ano.activated.connect(self.reset_combo_slides)
        self.combo_ano.activated.connect(self.reset_combo)
        self.combo_box_meses.activated.connect(self.reset_combo_slides)
        self.combo_box_slides.activated.connect(self.reset_combo)

        self.combo_box_layout.addWidget(self.combo_ano)
        self.combo_box_layout.addWidget(self.combo_box_meses)
        self.combo_box_layout.addWidget(self.combo_box_slides)
        
        self.layout.addLayout(self.combo_box_layout)


    def reset_combo_slides(self):
        self.combo_box_slides.setCurrentIndex(0)
    
    def reset_combo(self):
        self.combo_box_meses.setCurrentIndex(0)
    
    def criar_tabela(self):
        self.tabela = QtWidgets.QTableWidget()
        self.adicionar_funcionalidade_excluir(self.tabela)
        self.adicionar_funcionalidade_ordenar_colunas(self.tabela)
        self.tabela.setStyleSheet("""
            QTableWidget {
                gridline-color: #D3D3D3;
                background-color: #F4F4F4;
                alternate-background-color: #E8E8E8;
                border: 1px solid #A9A9A9;
                font-size: 14px;
            }
            
            QHeaderView::section {
                background-color: #001F3F;
                color: white;
                padding: 8px;
                border: 1px solid #A9A9A9;
                font-weight: bold;
                text-align: center;
            }
            
            QTableWidget::item {
                padding: 8px;
                border: 1px solid #D3D3D3;
                text-align: center;
                color: black;
            }
            
            QTableWidget::item:selected {
                background-color: #ADD8E6;
                color: black;
            }
        """)
        self.tabela.setAlternatingRowColors(True)
        self.layout.addWidget(self.tabela)
        

    def limpar_tabela(self, tabela):
        tabela.clearContents()
        tabela.setRowCount(0)
        tabela.setColumnCount(0)

    def adicionar_funcionalidade_excluir(self, tabela):
        def excluir():
            modelo_selecao = tabela.selectionModel()
            linhas_selecionadas = sorted(set(index.row() for index in modelo_selecao.selectedRows()))
            colunas_selecionadas = sorted(set(index.column() for index in modelo_selecao.selectedColumns()))

            if linhas_selecionadas:
                # Confirmação para excluir linhas
                resposta = QMessageBox.question(
                    tabela,
                    "Confirmação",
                    "Deseja realmente excluir as linhas selecionadas?" if len(linhas_selecionadas) > 1 else "Deseja realmente excluir a linha selecionada?",
                    QMessageBox.Yes | QMessageBox.No
                )
                if resposta == QMessageBox.Yes:
                    # Remove as linhas na ordem inversa para evitar problemas de índice
                    for linha in reversed(linhas_selecionadas):
                        tabela.removeRow(linha)

            elif colunas_selecionadas:
                # Confirmação para excluir colunas
                resposta = QMessageBox.question(
                    tabela,
                    "Confirmação",
                    "Deseja realmente excluir as colunas selecionadas?" if len(colunas_selecionadas) > 1 else "Deseja realmente excluir a coluna selecionada?",
                    QMessageBox.Yes | QMessageBox.No
                )
                if resposta == QMessageBox.Yes:
                    for coluna in reversed(colunas_selecionadas):
                        tabela.removeColumn(coluna)

            else:
                QMessageBox.warning(
                    tabela,
                    "Atenção",
                    "Selecione ao menos uma linha ou coluna para exclusão."
                )

    # Conecta o atalho Delete à função de exclusão
        excluir_atalho = QShortcut(QKeySequence("Delete"), tabela)
        excluir_atalho.activated.connect(excluir)
        excluir_atalho_backspace = QShortcut(QKeySequence("Backspace"), tabela)
        excluir_atalho_backspace.activated.connect(excluir)

    def adicionar_funcionalidade_ordenar_colunas(self, tabela):
        """
        Adiciona funcionalidade para alternar entre ordem crescente e decrescente ao clicar no cabeçalho da coluna.
        """
        # Dicionário para rastrear o estado da ordenação de cada coluna
        estados_ordenacao = {}

        def alternar_ordenacao(indice_coluna):
            """
            Alterna a ordenação entre crescente e decrescente para a coluna clicada.
            Realiza verificação para não ordenar colunas com letras (texto).
            """
            # Verifica se a coluna contém apenas números
            for row in range(tabela.rowCount()):
                item = tabela.item(row, indice_coluna)
                if item is not None:
                    try:
                        # Tenta converter o valor da célula para número
                        float(item.text().replace(",", "").replace(".", ""))
                    except ValueError:
                        return

            # Alterna o estado da ordenação
            ordem_atual = estados_ordenacao.get(indice_coluna, Qt.AscendingOrder)
            nova_ordem = Qt.DescendingOrder if ordem_atual == Qt.AscendingOrder else Qt.AscendingOrder
            estados_ordenacao[indice_coluna] = nova_ordem

            # Preserva o alinhamento das células ao substituir os itens
            for row in range(tabela.rowCount()):
                item = tabela.item(row, indice_coluna)
                if item is not None:
                    # Guarda o alinhamento original
                    alinhamento_original = item.textAlignment()

                    # Substitui o item por um NumericTableWidgetItem
                    tabela.setItem(row, indice_coluna, NumericTableWidgetItem(item.text()))

                    # Restaura o alinhamento original
                    tabela.item(row, indice_coluna).setTextAlignment(alinhamento_original)

            # Ordena os itens da tabela
            tabela.sortItems(indice_coluna, nova_ordem)

        # Conecta o evento de clique no cabeçalho da tabela à função de alternância
        tabela.horizontalHeader().sectionClicked.connect(alternar_ordenacao)        
    
    def gerar_tabela_se_valido(self, tabela, combo_text, empresa_selecionada, ano, stack_layout):
        if combo_text == "Slides" or combo_text == "":
            return

        if ano == "Escolha o ano" and combo_text not in ["Slides", "Balancetes"]:
            mensagem_error("Selecione um ano válido.")
            return
        if combo_text not in ["Slides", "Balancetes"]:
            self.gerar_tabela(tabela, combo_text, empresa_selecionada, ano, stack_layout)

    def gerar_tabela(self, tabela, tabela_selecionada, empresa, ano, stack_layout):
        dados_balancete, colunas = self.obter_dados_balancete(tabela_selecionada, empresa, ano)
        
        if dados_balancete is None or len(dados_balancete) == 0:
            mensagem_aviso("Nenhum dado encontrado para os critérios selecionados.")
            # msg_box = QtWidgets.QMessageBox()
            # msg_box.setIcon(QtWidgets.QMessageBox.Information)
            # msg_box.setWindowTitle("Informação")
            # msg_box.setText("Nenhum dado encontrado para os critérios selecionados.")
            # msg_box.setStandardButtons(QtWidgets.QMessageBox.Ok)
            # usar_icone(msg_box)
            stack_layout.setCurrentWidget(self.imagem_placeholder)
            # tabela.clearContents()
            # tabela.setRowCount(0)

            # msg_box.exec()
            return
        
        self.preencher_tabela_com_dados(tabela, dados_balancete, colunas, stack_layout)
    
    def obter_dados_balancete(self, tabela, empresa, ano):
        try:
            meses_dict = {
                "Janeiro": 1,
                "Fevereiro": 2,
                "Março": 3,
                "Abril": 4,
                "Maio": 5,
                "Junho": 6,
                "Julho": 7,
                "Agosto": 8,
                "Setembro": 9,
                "Outubro": 10,
                "Novembro": 11,
                "Dezembro": 12
            }
            
            with conectar_com_banco() as conn:
                with conn.cursor() as cursor:
                    meses = ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho', 'Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro'] 
                    cursor.execute("""
                        SELECT DISTINCT periodo FROM balancete_mensal WHERE empresa = %s AND periodo LIKE %s
                    """, (empresa, f"%{ano}%"))
                    periodos = [row[0] for row in cursor.fetchall()]  # Extrai os períodos únicos

                    if not periodos:
                        return [], []

                    # Gerar dinamicamente o select_part com base nos períodos disponíveis
                    select_part = ", ".join([
                        f"MAX(CASE WHEN periodo = '{periodo}' THEN COALESCE(`Saldo Atual`, 0) END) AS {meses[int(periodo.split('-')[0]) - 1]}"
                        for periodo in sorted(periodos)
                    ])

                    # Adicionar lógica para a coluna Total
                    ultimo_periodo = sorted(periodos)[-1]  # Último período disponível
                    soma_meses = " + ".join([f"COALESCE(MAX(CASE WHEN periodo = '{periodo}' THEN `Saldo Atual` END), 0)" for periodo in sorted(periodos)])

                    select_part += f""",
                        CASE
                            WHEN LEFT(Conta, 1) IN ('1', '2') THEN COALESCE(MAX(CASE WHEN periodo = '{ultimo_periodo}' THEN `Saldo Atual` END), 0)
                            WHEN LEFT(Conta, 1) = '3' THEN {soma_meses}
                        END AS Total
                    """

                    if tabela in meses_dict:
                        periodo = meses_dict[tabela]
                        query = f"""
                        WITH DistinctRows AS (
                            SELECT 
                                *, 
                                ROW_NUMBER() OVER (PARTITION BY Conta, `Descrição`, periodo ORDER BY `Data de Envio` DESC) AS row_num
                            FROM balancete_mensal
                            WHERE periodo = %s AND empresa = %s AND periodo LIKE %s
                        )
                        SELECT *
                        FROM DistinctRows
                        WHERE row_num = 1;
                        """
                        cursor.execute(query, (periodo, empresa, f"%{ano}%"))
                        dados = cursor.fetchall()
                        colunas = [desc[0] for desc in cursor.description]
                        return dados, colunas

                    elif tabela == "Balancete Acumulado":
                        query = f"""
                        WITH DistinctRows AS (
                            SELECT 
                                *, 
                                ROW_NUMBER() OVER (PARTITION BY Conta, `Descrição`, periodo ORDER BY `Data de Envio` DESC) AS row_num
                            FROM balancete_mensal
                            WHERE empresa = %s AND periodo LIKE %s
                        )
                        SELECT
                            Conta,
                            `Descrição`,
                            {select_part}
                        FROM
                            DistinctRows
                        WHERE
                            row_num = 1
                        GROUP BY
                            Conta,
                            `Descrição`
                        ORDER BY Conta;
                        """
                        cursor.execute(query, (empresa, f"%{ano}%"))
                        dados = cursor.fetchall()
                        colunas = [desc[0] for desc in cursor.description]
                        return dados, colunas

                    elif tabela not in meses_dict and tabela != "Balancete Acumulado":
                        cursor.execute("SELECT id from empresas where razao_social = %s", (empresa,))
                        id_empresa = cursor.fetchone()[0]
                        cursor.execute("""
                            SELECT id from slides where nome = %s and id_empresa = %s
                            """, (tabela, id_empresa))
                        id_slide = cursor.fetchone()[0]
                        cursor.execute("""
                            SELECT codigo FROM slide_codigos WHERE slide_id = %s
                        """, (id_slide,))
                        codigos = cursor.fetchall()

                        # Cria uma lista para armazenar as condições de LIKE
                        like_conditions = []
                        params = []

                        # Itera pelos códigos e adiciona cada condição à lista
                        for codigo in codigos:
                            like_conditions.append("conta LIKE %s")
                            params.append(f"{codigo[0]}%")  # Adiciona o código com wildcards para o LIKE

                        # Constrói a consulta com todas as condições de LIKE usando OR
                        query = f"""
                            WITH DistinctRows AS (
                                SELECT 
                                    *, 
                                    ROW_NUMBER() OVER (PARTITION BY Conta, `Descrição`, periodo ORDER BY `Data de Envio` DESC) AS row_num
                                FROM balancete_mensal
                                WHERE {' OR '.join(like_conditions)} AND periodo LIKE %s
                            )
                            SELECT Conta, `Descrição`, {select_part} 
                            FROM DistinctRows
                            WHERE row_num = 1
                            GROUP BY Conta, `Descrição`
                            ORDER BY Conta;
                        """
                        params.append(f"%{ano}%")
                        cursor.execute(query, params)
                        dados = cursor.fetchall()
                        colunas = [desc[0] for desc in cursor.description]
                        return dados, colunas

                    else:
                        return [], []  # Caso não entre em nenhuma das condições

        except mysql.connector.errors.ProgrammingError as e:
            if e.errno == 1146:  # Código para "tabela não existe"
                print(f"A tabela selecionada ('{tabela}') não existe no banco de dados.")
            else:
                raise
        return [], []

    def preencher_tabela_com_dados(self, tabela, dados, colunas, stack_layout):
        msg_box = QtWidgets.QMessageBox()
        msg_box.setIcon(QtWidgets.QMessageBox.Information)
        msg_box.setWindowTitle("Informação")
        msg_box.setText("Nenhum dado encontrado. A tabela selecionada não está disponível.")
        msg_box.setStandardButtons(QtWidgets.QMessageBox.Ok)
        
        if not dados:
            usar_icone(msg_box)
            # tabela.clearContents()
            # tabela.setRowCount(0)
            msg_box.exec()
            imagem_placeholder = QtWidgets.QLabel()
            imagem_placeholder.setPixmap(QtGui.QPixmap("images\\logo.png").scaled(300, 300, QtCore.Qt.KeepAspectRatio))
            imagem_placeholder.setAlignment(QtCore.Qt.AlignCenter)
            stack_layout.setCurrentWidget(imagem_placeholder)
            return
        
        
        if 'id' in colunas:
            id_index = colunas.index('id')
            colunas.pop(id_index)  # Remove o nome da coluna
            dados = [tuple(valor for i, valor in enumerate(linha) if i != id_index) for linha in dados]

        # Remover a coluna 'row_num' dos dados e das colunas
        if 'row_num' in colunas:
            row_num_index = colunas.index('row_num')
            colunas.pop(row_num_index)  # Remove o nome da coluna
            dados = [tuple(valor for i, valor in enumerate(linha) if i != row_num_index) for linha in dados]  # Remove valores correspondentes

        tabela.clearContents()
        tabela.setRowCount(0)
        tabela.setColumnCount(len(colunas))
        tabela.setHorizontalHeaderLabels(colunas)

        for linha_num, linha in enumerate(dados):
            tabela.insertRow(linha_num)
            for coluna_num, valor in enumerate(linha):
                # Formata os dados apenas para as colunas de meses: Janeiro, Fevereiro, Março, e outros campos numéricos
                if colunas[coluna_num] in ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho', 'Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro', 'Saldo Anterior', 'Débitos', 'Créditos', 'Saldo Atual', 'Total']:
        # Verifica se o valor não é None e tenta convertê-lo para float
                    if valor is not None:
                        try:
                            # Converte para float, remove parte decimal e formata como inteiro
                            valor_formatado = int(float(valor))  # Remove a parte após a vírgula completamente
                            # Adiciona separador de milhar
                            valor_formatado = "{:,}".format(valor_formatado).replace(",", ".")
                        except ValueError:
                            valor_formatado = str(valor)  # Se não for numérico, mantém o valor original
                    else:
                        valor_formatado = "0"  # Se o valor for None, define como "0"

                    # Cria o item com a formatação correta
                    item = QtWidgets.QTableWidgetItem(valor_formatado)
                    item.setTextAlignment(QtCore.Qt.AlignRight | QtCore.Qt.AlignVCenter)
                    item.setFlags(QtCore.Qt.ItemIsSelectable | QtCore.Qt.ItemIsEnabled)  # Alinha à direita e centraliza verticalmente
                else:
                    # Para outros tipos de dados, apenas converta sem separador
                    item = QtWidgets.QTableWidgetItem(str(valor))
                    item.setTextAlignment(QtCore.Qt.AlignVCenter | QtCore.Qt.AlignLeft)# Alinhamento centralizado por padrão
                    item.setFlags(QtCore.Qt.ItemIsSelectable | QtCore.Qt.ItemIsEnabled)
                tabela.setItem(linha_num, coluna_num, item)

        tabela.resizeColumnsToContents()
        stack_layout.setCurrentWidget(tabela)

        for coluna in range(tabela.columnCount()):
            largura_atual = tabela.columnWidth(coluna)
            tabela.setColumnWidth(coluna, largura_atual + 20)
    
    def criar_slide(self, combo_box, tabela, empresa, stack_layout):
        stack_layout.setCurrentWidget(self.imagem_placeholder)
        with conectar_com_banco() as conn:
            with conn.cursor() as cursor:
                cursor.execute("SELECT 1 FROM balancete_mensal WHERE empresa = %s LIMIT 1", (empresa,))
                tabela_vazia = cursor.fetchone() is None
    
    # Se a tabela estiver vazia, exibir uma mensagem e retornar
        if tabela_vazia:
            mensagem_aviso("Insira um balancete para continuar.")
            return
        
        # Cria uma janela de diálogo para inserir o nome do slide e os números das contas
        dialogo = QDialog(self)
        dialogo.setWindowTitle("Criar Novo Slide")
        
        dialogo.resize(800, 600) 
        
        layout = QVBoxLayout(dialogo)
        
        # Campo de entrada para o nome do slide
        nome_label = QLabel("Digite o nome do slide:", dialogo)
        nome_label.setStyleSheet("font-size: 16px;")
        layout.addWidget(nome_label)
        
        nome_input = QLineEdit(dialogo)
        nome_input.setStyleSheet("font-size: 16px; padding: 10px;")
        nome_input.setPlaceholderText("Nome do Slide")
        layout.addWidget(nome_input)
        
        # Consultar as contas disponíveis no banco de dados
        with conectar_com_banco() as conn:
            with conn.cursor() as cursor:
                cursor.execute("SELECT DISTINCT Conta, Descrição FROM balancete_mensal WHERE empresa = %s ORDER BY Conta", (empresa,))
                contas = cursor.fetchall()
        
        # Lista de contas no QListWidget
        contas_lista = QListWidget(dialogo)
        contas_lista.setSelectionMode(QListWidget.MultiSelection)  # Permite seleção múltipla
        
        # Adicionar estilo para melhorar a aparência
        contas_lista.setStyleSheet("""
            QListWidget {
                font-size: 14px;
                border: 1px solid #ccc;
                background-color: black;
                padding: 5px;
            }
            QListWidget::item {
                padding: 5px;
                border-radius: 5px;
            }
            QListWidget::item:selected {
                background-color: #4CAF50;  /* Cor de fundo para o item selecionado */
                color: white;
            }
            QListWidget::item:hover {
                background-color: #191970;  /* Cor ao passar o mouse */
            }
        """)
        contas_lista.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        # Adicionar as contas na lista
        for conta in contas:
            contas_lista.addItem(f"{conta[0]} - {conta[1]}")
        
        layout.addWidget(contas_lista)
        
        # Botão de salvar
        salvar_btn = QPushButton("Salvar", dialogo)
        salvar_btn.clicked.connect(lambda: self.salvar_numeros_contas(
            nome_input, contas_lista.selectedItems(), dialogo, combo_box, tabela, empresa))
        salvar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        salvar_btn.setStyleSheet("font-size: 16px; padding: 10px; background-color: #001F3F; color: white;")
        layout.addWidget(salvar_btn)
        
        dialogo.exec()
    
    def salvar_numeros_contas(self, nome_input, contas_selecionadas, dialogo, combo_box, tabela, empresa):
    # Obtém o nome do slide e os números das contas
        with conectar_com_banco() as conn:
            tabela_slide(conn)
            codigos_slide(conn)

        nome_slide = nome_input.text().strip()
        contas = [item.text().split(" - ")[0] for item in contas_selecionadas]  # Extrai o código da conta
        
        if not nome_slide:
            mensagem_aviso("Digite o nome do slide.")
            return
        elif not contas:
            mensagem_aviso("Selecione ao menos uma conta.")
            # print("Erro: Nome do slide ou contas inválidas.")
            return
        # Conexão com o banco de dados
        conexao = conectar_com_banco()
        cursor = conexao.cursor()

        try:
            # Inserir o slide na tabela 'slides'
            cursor.execute(
                "INSERT INTO slides (nome, id_empresa) VALUES (%s, (SELECT id FROM empresas WHERE razao_social = %s))",
                (nome_slide, empresa)
            )
            slide_id = cursor.lastrowid  

        
            for conta in contas:
                cursor.execute(
                    "INSERT INTO slide_codigos (slide_id, codigo) VALUES (%s, %s)",
                    (slide_id, conta)
                )

            conexao.commit()  # Confirma as alterações no banco de dados
            
            # Atualiza o combo_box com o novo slide
            combo_box.addItem(nome_slide)
            
            mensagem_sucesso("Slide criado com sucesso!")
            # Fecha o diálogo
            dialogo.accept()

        except Exception as e:
            # print(f"Erro ao salvar o slide: {e}")
            conexao.rollback()  # Desfaz alterações em caso de erro
        
        finally:
            cursor.close()
            conexao.close()


    def editar_slide(self, combo_box, tabela, empresa, stack_layout):
        stack_layout.setCurrentWidget(self.imagem_placeholder)
        with conectar_com_banco() as conexao:
            with conexao.cursor() as cursor:
                cursor.execute("SELECT id FROM empresas WHERE razao_social = %s", (empresa,))
                id_empresa = cursor.fetchone()[0]
                cursor.execute("SELECT id, nome FROM slides WHERE id_empresa = %s", (id_empresa,))
                slides = cursor.fetchall()
                cursor.execute("SELECT DISTINCT Conta, Descrição FROM balancete_mensal WHERE empresa = %s ORDER BY Conta", (empresa,))
                contas = cursor.fetchall()

        if not slides:
            mensagem_aviso("Não há slides cadastrados para a empresa selecionada.")
            return
        
        dialogo = QDialog(self)
        dialogo.setWindowTitle("Editar Slide")
        dialogo.resize(800, 600)

        layout = QVBoxLayout(dialogo)

        slide_label = QLabel("Selecione o slide a ser editado:", dialogo)
        layout.addWidget(slide_label)

        slide_combo = QComboBox(dialogo)
        slide_combo.setStyleSheet("font-size: 16px; padding: 10px;")

        
        # Associar o slide_id com o nome
        for slide in slides:
            slide_combo.addItem(slide[1], slide[0])  # Adiciona o nome e armazena o id como dado
        layout.addWidget(slide_combo)

        nome_label = QLabel("Digite o novo nome do slide:", dialogo)
        nome_label.setStyleSheet("font-size: 16px;")
        layout.addWidget(nome_label)

        nome_input = QLineEdit(dialogo)
        nome_input.setStyleSheet("font-size: 16px; padding: 5px;")
        layout.addWidget(nome_input)

        contas_label = QLabel("Selecione as contas do slide:", dialogo)
        layout.addWidget(contas_label)

        contas_lista = QListWidget(dialogo)
        contas_lista.setSelectionMode(QListWidget.MultiSelection)

        contas_lista.setStyleSheet(""" 
            QListWidget {
                font-size: 14px;
                border: 1px solid #ccc;
                background-color: black;
                padding: 5px;
            }
            QListWidget::item {
                padding: 5px;
                border-radius: 5px;
            }
            QListWidget::item:selected {
                background-color: #4CAF50;  /* Cor de fundo para o item selecionado */
                color: white;
            }
            QListWidget::item:hover {
                background-color: #191970;  /* Cor ao passar o mouse */
            }
        """)
        contas_lista.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))

        # Adicionar as contas na lista
        for conta in contas:
            contas_lista.addItem(f"{conta[0]} - {conta[1]}")
        
        layout.addWidget(contas_lista)

        salvar_btn = QPushButton("Salvar", dialogo)
        salvar_btn.clicked.connect(lambda: self.salvar_edicao_slide(
            slide_combo, nome_input, contas_lista, dialogo, combo_box, tabela, empresa))
        salvar_btn.setCursor(QtGui.QCursor(QtCore.Qt.PointingHandCursor))
        salvar_btn.setStyleSheet("font-size: 16px; padding: 10px; background-color: #001F3F; color: white;")
        layout.addWidget(salvar_btn)

        dialogo.exec()
    
    def salvar_edicao_slide(self, slide_combo, nome_input, contas_lista, dialogo, combo_box, tabela, empresa):
    # Primeiro, obter o ID da empresa com base no nome
        conexao = conectar_com_banco()
        cursor = conexao.cursor()
        
        try:
            cursor.execute("SELECT id FROM empresas WHERE razao_social = %s", (empresa,))
            id_empresa = cursor.fetchone()
            
            if id_empresa is None:
                # print("Erro: Empresa não encontrada.")
                return
            
            id_empresa = id_empresa[0]  # Agora temos o ID da empresa

            slide_id = slide_combo.currentData()  
            novo_nome = nome_input.text().strip()
            # print(empresa)
            contas_selecionadas = [item.text().split(" - ")[0] for item in contas_lista.selectedItems()]

            if not novo_nome:
                mensagem_aviso("Digite o nome do slide.")
                # print("Erro: Nome do slide inválido.")
                return
            elif not contas_selecionadas:
                mensagem_aviso("Selecione ao menos uma conta.")
                # print("Erro: Nenhuma conta selecionada.")
                return
            
            # Atualizar o nome do slide
            cursor.execute("UPDATE slides SET nome = %s WHERE id = %s", (novo_nome, slide_id))

            # Remover as contas antigas associadas ao slide
            cursor.execute("DELETE FROM slide_codigos WHERE slide_id = %s", (slide_id,))

            # Adicionar as contas selecionadas ao slide
            for conta_codigo in contas_selecionadas:
                cursor.execute("INSERT INTO slide_codigos (slide_id, codigo) VALUES (%s, %s)", (slide_id, conta_codigo))

            conexao.commit()

            mensagem_sucesso("Slide editado com sucesso!")

            # Atualizar a lista de slides no combo_box (tela principal)
            combo_box.clear()
            cursor.execute("SELECT id, nome FROM slides WHERE id_empresa = %s", (id_empresa,))
            slides_atualizados = cursor.fetchall()
            # print(f"Slides recuperados do banco: {slides_atualizados}")  # Verificar o retorno
            # if not slides_atualizados:
            #     print("Nenhum slide encontrado para a empresa.")
            
            combo_box.addItem("Slides")
            combo_box.model().item(0).setEnabled(False) 
            # Define o item "Slides" como nulo
            
            for slide in slides_atualizados:
                # A tupla de cada slide tem dois elementos, o id e o nome.
                combo_box.addItem(slide[1], slide[0])  # slide[1] = nome, slide[0] = id
                # print(f"Adicionado ao combo: {slide[1]}")  # Verificar se os itens estão sendo adicionados

            dialogo.accept()

        except Exception as e:
            mensagem_error(f"Erro ao editar o slide: {e}")
            conexao.rollback()
        
        finally:
            cursor.close()
            conexao.close()

    def exportar_tabela_para_powerpoint(self, tabela, nome_slide, nome_empresa, caminho_logo):

        if self.stack_layout.currentWidget() == self.imagem_placeholder:
            mensagem_error("Nenhum dado encontrado para exportar.")
            return
    # Obter o caminho da pasta Downloads
        caminho_downloads = Path.home() / "Downloads" / "Apresentações" / nome_empresa
        caminho_downloads.mkdir(parents=True, exist_ok=True)
        arquivo_saida = caminho_downloads / f"{nome_empresa} - Apresentação.pptx"

        # Verificar se a tabela tem dados válidos para exportar
        linhas = tabela.rowCount()
        colunas = tabela.columnCount()
        if linhas == 0 or colunas == 0:
            mensagem_error("Nenhum dado encontrado para exportar.")
            return

        # Tentar abrir uma apresentação existente, caso contrário, criar uma nova
        try:
            prs = Presentation(arquivo_saida)  # Tenta carregar uma apresentação existente
        except Exception:
            prs = Presentation()  # Se o arquivo não existir, cria uma nova apresentação

        # Adiciona um novo slide com layout de título + conteúdo
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Adiciona título ao slide
        slide.shapes.title.text = f"{nome_slide}"
        titulo = slide.shapes.title
        titulo.text_frame.paragraphs[0].font.name = "Arial"
        titulo.text_frame.paragraphs[0].font.bold = True
        titulo.text_frame.paragraphs[0].font.size = Pt(18)
        titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        titulo.top = Inches(0.25)
        titulo.width = Inches(9.5)
        titulo.height = Inches(1)

        # Adiciona a logo da empresa ao lado do título
        try:
            slide.shapes.add_picture(
                caminho_logo,
                Inches(7),
                Inches(0.25),
                width=Inches(2.5),
                height=Inches(1)
            )
        except Exception as e:
            print(f"Erro ao adicionar a imagem da logo: {e}")

        # Configuração do fundo
        fundo = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.4), Inches(10), Inches(6.6)
        )
        fundo.fill.solid()
        fundo.fill.fore_color.rgb = RGBColor(240, 240, 240)
        fundo.line.fill.solid()
        fundo.line.fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Recupera os dados da tabela
        dados = []
        for i in range(linhas):
            row_dados = []
            for j in range(colunas):
                item = tabela.item(i, j)
                if item:
                    row_dados.append(item.text())
                else:
                    row_dados.append("")
            dados.append(row_dados)

        # Identificar o índice da coluna "Descrição"
        coluna_descricao = -1
        for col_index in range(colunas):
            if tabela.horizontalHeaderItem(col_index).text() == "Descrição":
                coluna_descricao = col_index
                break

        # Insere os dados como uma tabela no slide
        tabela_slide = slide.shapes.add_table(
            rows=len(dados) + 1,
            cols=colunas,
            left=Inches(1.4),
            top=Inches(1.6),
            width=Inches(4),
            height=Inches(0.5),
        ).table

        # Configurar os cabeçalhos da tabela
        for col_index in range(colunas):
            tabela_slide.cell(0, col_index).text = tabela.horizontalHeaderItem(col_index).text()
            tabela_slide.cell(0, col_index).text_frame.paragraphs[0].font.bold = True
            tabela_slide.cell(0, col_index).text_frame.paragraphs[0].font.size = Pt(10)

        colunas_direita = ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho', 'Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro', 'Saldo Anterior', 'Débitos', 'Créditos', 'Saldo Atual', 'Total']

        for row_index, linha_dados in enumerate(dados):
            for col_index, valor in enumerate(linha_dados):
                cell = tabela_slide.cell(row_index + 1, col_index)
                cell.text = valor
                cell.text_frame.word_wrap = False
                cell.text_frame.paragraphs[0].font.size = Pt(8)

                if tabela.horizontalHeaderItem(col_index).text() in colunas_direita:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                else:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        if coluna_descricao != -1:
            tabela_slide.columns[coluna_descricao].width = Inches(3.5)

        for col_index in range(colunas):
            if col_index != coluna_descricao:
                tabela_slide.columns[col_index].width = Inches(0.9)

        altura_total = 0.4 + 0.3 * len(dados)
        tabela_slide.height = Inches(altura_total)

        try:
            prs.save(arquivo_saida)
        except PermissionError:
            mensagem_error("Erro ao salvar a apresentação. Verifique se o arquivo não está aberto.")
            return
        except Exception as e:
            mensagem_error("Erro ao salvar a apresentação.")
            return

        mensagem_sucesso(f"Apresentação exportada com sucesso para o arquivo '{arquivo_saida}'.")

    def baixar_icone(self, url, caminho):
        dir_name = os.path.dirname(caminho)
        if not os.path.exists(dir_name):
            os.makedirs(dir_name)
        
        resposta = requests.get(url)
        with open(caminho, 'wb') as arquivo:
            arquivo.write(resposta.content)
    
    def recurso_caminho(self, relativo):
        if hasattr(sys, '_MEIPASS'): 
            return os.path.join(sys._MEIPASS, relativo)
        return os.path.join(os.path.abspath("."), relativo)

class NumericTableWidgetItem(QTableWidgetItem):
    def __lt__(self, other):
        """
        Sobrescreve o operador de comparação para lidar com valores numéricos.
        """
        try:
            # Converte os textos para números para comparar
            return float(self.text().replace(",", "").replace(".", "")) < float(other.text().replace(",", "").replace(".", ""))
        except ValueError:
            # Se não for numérico, usa a comparação padrão de strings
            return super().__lt__(other)

def main():
    app = QtWidgets.QApplication(sys.argv)
    login = UserLogin()
    usar_icone(login)
    login.showMaximized()
    sys.exit(app.exec())

if __name__ == '__main__':
    main()